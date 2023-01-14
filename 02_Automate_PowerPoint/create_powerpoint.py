import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Determine the directory of the script
script_dir = os.path.dirname(os.path.abspath(__file__))

input_folder = os.path.join(script_dir, 'input')
charts_folder = os.path.join(script_dir, 'charts')
ppt_file = 'financial_data.pptx'

# Create the charts folder if it doesn't exist
if not os.path.exists(charts_folder):
    os.mkdir(charts_folder)

# Create a new PowerPoint presentation
prs = Presentation()

# Iterate through all Excel files in the input folder
for excel_file in os.listdir(input_folder):
    if not excel_file.endswith('.xlsx'):
        continue

    # Read the financial data from the first worksheet of the Excel file
    file_path = os.path.join(input_folder, excel_file)
    df = pd.read_excel(file_path, sheet_name=0, usecols="A:P")
    df = df.dropna()

    # Group the data by the "Product" column and sum up the "Sales" column
    grouped = df.groupby('Product').sum()['Sales']

    # Create a chart using the seaborn library
    sns.barplot(x=grouped.index, y=grouped.values)
    plt.title(excel_file)
    plt.xlabel('Product')
    plt.ylabel('Sales')
    plt.tight_layout()

    # Save the chart to the charts folder
    chart_file = excel_file.replace('.xlsx', '.png')
    chart_path = os.path.join(charts_folder, chart_file)
    plt.savefig(chart_path)

    # Add a slide to the PowerPoint presentation and insert the chart and title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = excel_file.replace('.xlsx','')

    chart_file = chart_path
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6)
    slide.shapes.add_picture(chart_file, left, top, width=width, height=height)

# Save the PowerPoint presentation in the same directory as the script
ppt_path = os.path.join(script_dir, ppt_file)
prs.save(ppt_path)
